import re
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import requests

def fetch_unsplash_image(query, filename="temp.jpg"):
    url = f"https://source.unsplash.com/800x600/?{query}"
    img_data = requests.get(url).content
    with open(filename, "wb") as f:
        f.write(img_data)
    return filename

def set_slide_background(slide, rgb_color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color

def create_chart(data_dict, filename="chart.png"):
    """Generate a bar chart from a dict and save as image"""
    plt.figure(figsize=(4,3))
    plt.bar(data_dict.keys(), data_dict.values(), color='skyblue')
    plt.tight_layout()
    plt.savefig(filename)
    plt.close()
    return filename

def parse_numbers(lines):
    """Detect simple data patterns like 'Sales 30' and return dict"""
    data = {}
    for line in lines:
        match = re.match(r"(.+?)\s+(\d+)", line)
        if match:
            key, value = match.groups()
            data[key] = int(value)
    return data if data else None

def script_to_ppt(script_text, theme="corporate", output_file="chart_slides.pptx"):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    lines = script_text.strip().split("\n")
    title, content_lines = None, []

    themes = {
        "corporate": [RGBColor(230,245,255), RGBColor(220,240,255), RGBColor(245,245,245)],
        "modern": [RGBColor(255,240,230), RGBColor(235,255,235), RGBColor(245,235,255)],
        "dark": [RGBColor(30,30,30), RGBColor(50,50,50), RGBColor(20,20,40)]
    }
    palette = themes.get(theme, themes["corporate"])
    color_index = 0

    def add_slide(title, content_lines, color_index):
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, palette[color_index % len(palette)])

        # Title
        t = slide.shapes.title
        t.text = title
        f = t.text_frame.paragraphs[0].font
        f.size = Pt(36)
        f.bold = True
        f.color.rgb = RGBColor(0,70,140) if theme!="dark" else RGBColor(255,255,255)

        # Content
        c = slide.placeholders[1]
        tf = c.text_frame
        tf.clear()
        for line in content_lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(50,50,50) if theme!="dark" else RGBColor(220,220,220)
            p.level = 1 if line.startswith("-") else 0

        # Insert chart if numbers detected
        data = parse_numbers(content_lines)
        if data:
            chart_file = create_chart(data)
            slide.shapes.add_picture(chart_file, Inches(5), Inches(1.5), Inches(3.5), Inches(2.5))
        else:
            # Otherwise insert stock image
            try:
                img_file = fetch_unsplash_image(title)
                slide.shapes.add_picture(img_file, Inches(5), Inches(1.5), Inches(3.5), Inches(2.5))
            except Exception as e:
                print(f"⚠️ Could not fetch image for {title}: {e}")

    for line in lines:
        if line.startswith("# "):
            if title:
                add_slide(title, content_lines, color_index)
                color_index += 1
            title = line[2:].strip()
            content_lines = []
        else:
            if line.strip():
                content_lines.append(line.strip())

    if title:
        add_slide(title, content_lines, color_index)

    prs.save(output_file)
    print(f"✅ Slides with charts saved to {output_file}")


# Example usage
script_text = """
# Sales Overview
Product A 30
Product B 45
Product C 25

# Problem
- Manual work is slow
- Data is overwhelming

# Solution
AI automation improved efficiency by 50
"""
script_to_ppt(script_text, theme="modern", output_file="sales_slides.pptx")
